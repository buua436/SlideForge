from __future__ import annotations

import hashlib
import tempfile
from pathlib import Path
from typing import Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from PIL import Image
import lxml.etree as ET

from ppt_ui.core.component import RenderContext
from ppt_ui.core.diagnostics import Diagnostic
from ppt_ui.core.layout import Box, layout_from_spec, resolve_block_box
from ppt_ui.core.text_overflow import apply_text_overflow
from ppt_ui.icons.provider import IconRegistry, IconRequest, default_icon_registry
from ppt_ui.primitives import (
    ChartPrimitive,
    Connector,
    Ellipse,
    Group,
    IconPrimitive,
    ImagePrimitive,
    Line as LinePrimitive,
    PathPrimitive,
    Polygon,
    Rect,
    RichText,
    TablePrimitive,
    Text,
)
from ppt_ui.styles import Style, StyleResolver, StyleSheet, StyleTarget


def _rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _set_shape_opacity(shape: object, opacity: float) -> None:
    sp = shape._element
    spPr = sp.spPr
    solidFill = spPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = ET.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', str(int(opacity * 100000)))


def _add_native_shadow(shape: object, blur_radius: int, distance: int, direction: int = 5400000, opacity: float = 0.4) -> None:
    spPr = shape._element.spPr
    effectLst = spPr.find(qn('a:effectLst'))
    if effectLst is None:
        effectLst = ET.SubElement(spPr, qn('a:effectLst'))
    outerShdw = ET.SubElement(effectLst, qn('a:outerShdw'))
    outerShdw.set('blurRad', str(blur_radius))
    outerShdw.set('dist', str(distance))
    outerShdw.set('dir', str(direction))
    outerShdw.set('algn', 'tl')
    srgbClr = ET.SubElement(outerShdw, qn('a:srgbClr'))
    srgbClr.set('val', '000000')
    alpha = ET.SubElement(srgbClr, qn('a:alpha'))
    alpha.set('val', str(int(opacity * 100000)))


def _inheritable_style(style: Style | None) -> Style:
    if style is None:
        return Style()
    return Style(
        color=style.color,
        font_family=style.font_family,
        font_size=style.font_size,
        font_weight=style.font_weight,
        align=style.align,
        valign=style.valign,
        line_spacing=style.line_spacing,
    )


def _merge_inherited_style(parent: Style | None, child: Style | None) -> Style:
    return _inheritable_style(parent).merge(_inheritable_style(child))


class PptxRenderer:
    def __init__(self, theme: object, icon_registry: IconRegistry | None = None) -> None:
        self.theme = theme
        self.prs = Presentation()
        self.prs.slide_width = Inches(theme.slide_width)
        self.prs.slide_height = Inches(theme.slide_height)
        self.icon_registry = icon_registry or default_icon_registry()
        self.icon_cache_dir = Path(tempfile.gettempdir()) / "slideforge" / "icons"
        self.diagnostics: list[Diagnostic] = []

    def render_deck(self, deck: object) -> None:
        visible_pages = [page for page in deck.pages if not page.hidden]
        total_pages = len(visible_pages)
        for page_index, page in enumerate(visible_pages, start=1):
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            ctx = RenderContext(slide=slide, theme=self.theme, renderer=self, stylesheet=deck.styles, component_registry=deck.components)
            self.render_page(ctx, deck, page, page_index, total_pages)

    def render_page(self, ctx: RenderContext, deck: object, page: object, page_index: int, total_pages: int) -> None:
        master = None
        if page.use_master:
            master_name = page.master or deck.default_master
            master = deck.masters.get(master_name)

        layout = layout_from_spec(page.layout, self.theme)
        if master is not None:
            master.render_background(ctx, page)
            master.render_back_chrome(ctx, page, layout, page_index, total_pages)
        else:
            self.background(ctx.slide)

        for block in page.blocks:
            if not block.visible:
                continue
            component = deck.components.create(block.type, block.props, variant=block.variant)
            theme_style = self.theme.component_default_style(block.type, block.variant)
            component_style = {**theme_style, **block.style}
            component.render(
                ctx.with_block(block_id=block.id, class_names=block.class_names).with_style(component_style),
                resolve_block_box(layout, block.layout),
            )

        if master is not None:
            master.render_fore_chrome(ctx, page, layout, page_index, total_pages)

    def save(self, output_path: str | Path) -> Path:
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(path)
        return path

    def render_tree(
        self,
        slide: object,
        primitive: object,
        *,
        stylesheet: StyleSheet | None = None,
    ) -> list[object]:
        """Render a primitive tree to a PPT slide."""

        resolver = StyleResolver(self.theme, stylesheet or StyleSheet())
        return self.render_primitive(slide, primitive, resolver=resolver)

    def render_primitive(
        self,
        slide: object,
        primitive: object,
        *,
        resolver: StyleResolver | None = None,
        inherited_style: Style | None = None,
    ) -> list[object]:
        """Render one primitive node and its children."""

        resolver = resolver or StyleResolver(self.theme)
        shapes: list[object] = []
        if isinstance(primitive, Group):
            group_style = self._primitive_style(primitive, resolver, base=inherited_style)
            child_inherited = _merge_inherited_style(inherited_style, group_style)
            for child in primitive.children:
                shapes.extend(self.render_primitive(slide, child, resolver=resolver, inherited_style=child_inherited))
            return shapes

        style = self._primitive_style(primitive, resolver, base=inherited_style)
        rendered = self._render_single_primitive(slide, primitive, style)
        if isinstance(rendered, list):
            shapes.extend(item for item in rendered if item is not None)
        elif rendered is not None:
            shapes.append(rendered)

        for child in getattr(primitive, "children", ()):
            shapes.extend(self.render_primitive(slide, child, resolver=resolver, inherited_style=_inheritable_style(style)))
        return shapes

    def _primitive_style(self, primitive: object, resolver: StyleResolver, *, base: Style | None = None) -> Style:
        type_name = f"primitive.{getattr(primitive, 'type', 'primitive')}"
        target = StyleTarget(
            type_name=type_name,
            id=getattr(primitive, "id", None),
            class_names=tuple(getattr(primitive, "class_names", ())),
        )
        style = resolver.resolve(target, base=base, inline=getattr(primitive, "style", None))
        metadata = getattr(primitive, "metadata", {}) or {}
        component_type = metadata.get("component_type")
        slot_name = metadata.get("slot")
        if component_type and slot_name:
            slot_target = StyleTarget(
                type_name=str(component_type),
                id=getattr(primitive, "id", None),
                class_names=tuple(getattr(primitive, "class_names", ())),
                slot_name=str(slot_name),
            )
            style = resolver.resolve(slot_target, base=style)
        return style

    def _render_single_primitive(self, slide: object, primitive: object, style: Style) -> object | list[object] | None:
        box = getattr(primitive, "box", None)

        if isinstance(primitive, Text):
            if box is None:
                return None
            return self.text(
                slide,
                box,
                primitive.text,
                size=style.font_size,
                color=style.color or self.theme.colors.text_primary,
                bold=style.bold,
                align=style.align or "left",
                valign=style.valign or "top",
                font=style.font_family,
                line_spacing=style.line_spacing,
                overflow=str(style.extras.get("overflow", "fit")),
                max_chars=int(style.extras["max_chars"]) if style.extras.get("max_chars") is not None else None,
            )

        if isinstance(primitive, RichText):
            if box is None:
                return None
            text_value = "".join(run.text for run in primitive.runs)
            return self.text(
                slide,
                box,
                text_value,
                size=style.font_size,
                color=style.color or self.theme.colors.text_primary,
                bold=style.bold,
                align=style.align or "left",
                valign=style.valign or "top",
                font=style.font_family,
                line_spacing=style.line_spacing,
                overflow=str(style.extras.get("overflow", "fit")),
                max_chars=int(style.extras["max_chars"]) if style.extras.get("max_chars") is not None else None,
            )

        if isinstance(primitive, Rect):
            if box is None:
                return None
            if style.shadow:
                self.rect(
                    slide,
                    Box(box.x + self.theme.shadow.card_offset_x, box.y + self.theme.shadow.card_offset_y, box.w, box.h),
                    self.theme.colors.shadow_card,
                    rounded=style.radius if style.radius is not None else self.theme.radius_tokens.md,
                    opacity=0.65,
                )
            return self.rect(
                slide,
                box,
                style.fill or self.theme.colors.surface_white,
                line=style.stroke,
                rounded=style.radius if style.radius is not None else False,
                line_width=style.stroke_width or 0.55,
                opacity=style.opacity if style.opacity is not None else 1.0,
            )

        if isinstance(primitive, Ellipse):
            if box is None:
                return None
            return self.circle(
                slide,
                box,
                style.fill or self.theme.colors.surface_white,
                line=style.stroke,
                line_width=style.stroke_width or 0.55,
                opacity=style.opacity if style.opacity is not None else 1.0,
            )

        if isinstance(primitive, LinePrimitive):
            if box is not None and primitive.x1 == primitive.x2 == primitive.y1 == primitive.y2 == 0:
                return self.line(slide, box.x, box.y + box.h / 2, box.x + box.w, box.y + box.h / 2, color=style.stroke or style.color, width=style.stroke_width or 1.0)
            return self.line(slide, primitive.x1, primitive.y1, primitive.x2, primitive.y2, color=style.stroke or style.color, width=style.stroke_width or 1.0)

        if isinstance(primitive, ImagePrimitive):
            if box is None:
                return None
            return self.picture(slide, primitive.src, box, fit=primitive.fit)

        if isinstance(primitive, IconPrimitive):
            if box is None:
                return None
            props = primitive.icon_props or {}
            request = self.icon_registry.create_request(
                primitive.name,
                color=style.color or style.stroke or self.theme.colors.primary,
                size=int(props.get("size", 128)),
                width=int(props["width"]) if props.get("width") is not None else None,
                height=int(props["height"]) if props.get("height") is not None else None,
                rotate=props.get("rotate"),
                flip=str(props.get("flip", "")),
                stroke_width=float(props["stroke_width"]) if props.get("stroke_width") is not None else None,
            )
            if request is not None:
                shape = self.icon_picture(slide, request, box, opacity=style.opacity if style.opacity is not None else 1.0)
                if shape is not None:
                    return shape
            return self.text(slide, box, primitive.name[:2].upper(), size=style.font_size or 10, color=style.color or self.theme.colors.primary, bold=True, align="center", valign="middle")

        if isinstance(primitive, TablePrimitive):
            if box is None:
                return None
            self.add_table(slide, box, primitive.headers, primitive.rows)
            return []

        if isinstance(primitive, ChartPrimitive):
            if box is None:
                return None
            return self._render_chart_primitive(slide, primitive, box, style)

        if isinstance(primitive, Polygon):
            return self._render_polygon_primitive(slide, primitive, style)

        if isinstance(primitive, PathPrimitive):
            return None

        return None

    def _render_polygon_primitive(self, slide: object, primitive: Polygon, style: Style) -> list[object]:
        points = list(primitive.points)
        if len(points) < 2:
            return []
        shapes: list[object] = []
        for start, end in zip(points, [*points[1:], points[0]]):
            shape = self.line(slide, start[0], start[1], end[0], end[1], color=style.stroke or style.color, width=style.stroke_width or 1.0)
            if shape is not None:
                shapes.append(shape)
        return shapes

    def _render_chart_primitive(self, slide: object, primitive: ChartPrimitive, box: Box, style: Style) -> object | None:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        chart_type = primitive.chart_type.lower()
        data = CategoryChartData()
        palette = [color.lstrip("#") for color in self.theme.chart_palette]

        if chart_type in {"pie", "donut"}:
            labels = primitive.labels or tuple(f"Item {index + 1}" for index in range(len(primitive.values)))
            values = primitive.values or (1.0,)
            data.categories = labels[: len(values)]
            data.add_series("Share", values)
            ppt_chart_type = XL_CHART_TYPE.DOUGHNUT if chart_type == "donut" else XL_CHART_TYPE.PIE
        else:
            series = primitive.series
            value_count = max((len(item.values) for item in series), default=0)
            categories = primitive.categories or tuple(str(index + 1) for index in range(value_count))
            data.categories = categories
            for item in series:
                data.add_series(item.name, item.values)
            ppt_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED if chart_type == "bar" else XL_CHART_TYPE.LINE_MARKERS

        chart_shape = slide.shapes.add_chart(ppt_chart_type, Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h), data)
        chart = chart_shape.chart
        chart.has_title = False
        chart.has_legend = chart_type not in {"pie", "donut"} and len(primitive.series) > 1

        try:
            if chart_type in {"pie", "donut"}:
                for index, point in enumerate(chart.series[0].points):
                    point.format.fill.solid()
                    point.format.fill.fore_color.rgb = _rgb(palette[index % len(palette)])
            else:
                for index, series in enumerate(chart.series):
                    if chart_type == "line":
                        series.format.line.color.rgb = _rgb(palette[index % len(palette)])
                        series.format.line.width = Pt(style.stroke_width or 1.5)
                    else:
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = _rgb(palette[index % len(palette)])
        except (AttributeError, IndexError):
            pass

        return chart_shape

    def rgb(self, value: str) -> RGBColor:
        return _rgb(value)

    def background(self, slide: object) -> None:
        gradient = getattr(self.theme, "gradient", None)
        if gradient and gradient.stops:
            fill = slide.background.fill
            fill.gradient()
            fill.gradient_stops[0].color.rgb = _rgb(gradient.stops[0].color)
            fill.gradient_stops[1].color.rgb = _rgb(gradient.stops[-1].color)
            for i, stop in enumerate(gradient.stops):
                if i < len(fill.gradient_stops):
                    fill.gradient_stops[i].color.rgb = _rgb(stop.color)
                    fill.gradient_stops[i].position = stop.position / 100
            gf = fill._fill._gradFill
            lin = gf.find(qn('a:lin'))
            if lin is not None:
                lin.set('ang', str(gradient.angle))
                lin.set('scaled', '1')
        else:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(self.theme.colors.background)
        bg_pattern = getattr(self.theme, "background_pattern", "")
        if bg_pattern:
            self._apply_background_pattern(slide, bg_pattern)

    def _apply_background_pattern(self, slide: object, pattern: str) -> None:
        w = self.theme.slide_width
        h = self.theme.slide_height
        if pattern == "scanlines":
            for y_pos in [i * 0.15 for i in range(int(h / 0.15))]:
                self.rect(slide, Box(0, y_pos, w, 0.02), "000000", opacity=0.06)
        elif pattern == "grid":
            for x_pos in [i * 0.5 for i in range(int(w / 0.5) + 1)]:
                self.line(slide, x_pos, 0, x_pos, h, color="FFFFFF", width=0.3)
            for y_pos in [i * 0.5 for i in range(int(h / 0.5) + 1)]:
                self.line(slide, 0, y_pos, w, y_pos, color="FFFFFF", width=0.3)

    def rect(
        self,
        slide: object,
        box: Box,
        fill: str,
        line: str | None = None,
        *,
        rounded: bool | float = False,
        line_width: float = 0.55,
        opacity: float = 1.0,
        dash_style: int | None = None,
    ) -> object:
        if box.w <= 0 or box.h <= 0:
            return None
        use_rounded = bool(rounded)
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if use_rounded else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
        if use_rounded and len(shape.adjustments):
            radius = rounded if isinstance(rounded, (int, float)) and rounded is not True else self.theme.radius_tokens.md
            shape.adjustments[0] = radius
        if opacity < 1.0:
            _set_shape_opacity(shape, opacity)
        try:
            shape.shadow.inherit = False
        except AttributeError:
            pass
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = _rgb(line)
            shape.line.width = Pt(line_width)
            if dash_style is not None:
                shape.line.dash_style = dash_style
        return shape

    def line(
        self,
        slide: object,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        color: str | None = None,
        width: float = 1.0,
    ) -> object:
        if x1 == x2 and y1 == y2:
            return None
        shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        try:
            shape.shadow.inherit = False
        except AttributeError:
            pass
        shape.line.color.rgb = _rgb(color or self.theme.colors.primary)
        shape.line.width = Pt(width)
        return shape

    def circle(self, slide: object, box: Box, fill: str, line: str | None = None, line_width: float = 0.55, opacity: float = 1.0) -> object:
        if box.w <= 0 or box.h <= 0:
            return None
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
        if opacity < 1.0:
            _set_shape_opacity(shape, opacity)
        try:
            shape.shadow.inherit = False
        except AttributeError:
            pass
        if line is None:
            shape.line.fill.background()
        else:
            shape.line.color.rgb = _rgb(line)
            shape.line.width = Pt(line_width)
        return shape

    def picture(self, slide: object, src: str | Path, box: Box, *, fit: str = "contain") -> object:
        path = Path(src)
        if not path.exists() or not path.is_file() or box.w <= 0 or box.h <= 0:
            return None
        if fit == "stretch":
            return slide.shapes.add_picture(str(path), Inches(box.x), Inches(box.y), width=Inches(box.w), height=Inches(box.h))

        with Image.open(path) as image:
            image_w, image_h = image.size
        if image_w <= 0 or image_h <= 0:
            return None

        image_aspect = image_w / image_h
        box_aspect = box.w / box.h
        if fit == "cover":
            render_w = box.h * image_aspect if image_aspect < box_aspect else box.w
            render_h = box.w / image_aspect if image_aspect >= box_aspect else box.h
        else:
            render_w = box.w if image_aspect >= box_aspect else box.h * image_aspect
            render_h = box.w / image_aspect if image_aspect >= box_aspect else box.h
        x = box.x + (box.w - render_w) / 2
        y = box.y + (box.h - render_h) / 2
        return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(render_w), height=Inches(render_h))

    def icon_picture(self, slide: object, request: IconRequest, box: Box, *, opacity: float = 1.0) -> object:
        svg = self.icon_registry.resolve_svg(request)
        if svg is None:
            return None

        self.icon_cache_dir.mkdir(parents=True, exist_ok=True)
        cache_key = hashlib.sha1(f"{request}|{opacity}|{svg}".encode("utf-8")).hexdigest()
        png_path = self.icon_cache_dir / f"{cache_key}.png"
        if not png_path.exists():
            width = max(16, request.width or 128)
            height = max(16, request.height or request.width or 128)
            try:
                import resvg_py

                png_path.write_bytes(resvg_py.svg_to_bytes(svg_string=svg, width=width, height=height))
            except Exception:
                return None
            if opacity < 1.0:
                try:
                    with Image.open(png_path).convert("RGBA") as image:
                        alpha = image.getchannel("A").point(lambda value: int(value * max(0.0, min(1.0, opacity))))
                        image.putalpha(alpha)
                        image.save(png_path)
                except OSError:
                    return None
        return self.picture(slide, png_path, box, fit="contain")

    def text(
        self,
        slide: object,
        box: Box,
        text: str,
        size: int | None = None,
        color: str | None = None,
        bold: bool = False,
        align: str = "left",
        valign: str = "top",
        font: str | None = None,
        line_spacing: int | None = None,
        overflow: str = "fit",
        max_chars: int | None = None,
    ) -> object:
        if not text or box.w <= 0 or box.h <= 0:
            return None
        overflow_result = apply_text_overflow(text, overflow=overflow, max_chars=max_chars)
        if overflow_result.truncated:
            self.diagnostics.append(
                Diagnostic(
                    "warning",
                    "TEXT_TRUNCATED",
                    "Text was truncated by overflow policy.",
                    suggestion="Increase the box size, reduce the text length, or use overflow: fit.",
                )
            )
        shape = slide.shapes.add_textbox(Inches(box.x), Inches(box.y), Inches(box.w), Inches(box.h))
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        overflow_mode = overflow.lower().strip()
        frame.auto_size = {
            "clip": MSO_AUTO_SIZE.NONE,
            "resize": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
            "shape": MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT,
            "truncate": MSO_AUTO_SIZE.NONE,
        }.get(overflow_mode, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
        frame.margin_left = Inches(0)
        frame.margin_right = Inches(0)
        frame.margin_top = Inches(0)
        frame.margin_bottom = Inches(0)
        frame.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }.get(valign, MSO_ANCHOR.TOP)
        p = frame.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        if line_spacing is not None:
            pPr = p._pPr
            if pPr is None:
                pPr = ET.SubElement(p._p, qn('a:pPr'))
            lnSpc = ET.SubElement(pPr, qn('a:lnSpc'))
            spcPct = ET.SubElement(lnSpc, qn('a:spcPct'))
            spcPct.set('val', str(line_spacing))
        run = p.add_run()
        run.text = overflow_result.text
        latin = font or getattr(self.theme.fonts, "latin_font", "") or self.theme.fonts.family
        ea = self.theme.fonts.family
        run.font.name = latin
        rPr = run._r.get_or_add_rPr()
        ea_elem = rPr.find(qn('a:ea'))
        if ea_elem is not None:
            ea_elem.set('typeface', ea)
        else:
            ea_elem = ET.SubElement(rPr, qn('a:ea'))
            ea_elem.set('typeface', ea)
        run.font.size = Pt(size or self.theme.fonts.body_size)
        run.font.bold = bold
        run.font.color.rgb = _rgb(color or self.theme.colors.text_primary)
        return shape

    def bullet_list(
        self,
        slide: object,
        box: Box,
        items: Iterable[str],
        *,
        size: int | None = None,
        bullet_color: str | None = None,
        text_color: str | None = None,
    ) -> None:
        values = list(items)
        rows = box.split_rows(max(1, len(values)), gutter=0.04)
        for row, item in zip(rows, values):
            self.circle(slide, Box(row.x, row.y + 0.07, 0.075, 0.075), bullet_color or self.theme.colors.primary)
            self.text(
                slide,
                Box(row.x + 0.17, row.y, row.w - 0.17, row.h),
                item,
                size=size or self.theme.fonts.body_size,
                color=text_color or self.theme.colors.text_secondary,
            )

    def card(self, slide: object, box: Box, fill: str | None = None, line: str | None = None) -> object:
        return self.add_card(slide, box, fill=fill, line=line)

    def pill(self, slide: object, box: Box, text: str, fill: str | None = None, color: str | None = None) -> None:
        self.rect(slide, box, fill or self.theme.colors.primary, rounded=True)
        self.text(slide, box.inset(0.03, 0.0), text, size=self.theme.fonts.caption_size, color=color or "FFFFFF", bold=True, align="center", valign="middle")

    def accent_bar(self, slide: object, x: float = 0.28, y: float = 0.22, h: float = 0.72) -> None:
        self.add_accent_bar(slide, x=x, y=y, h=h)

    def content_box(self) -> Box:
        margin = self.theme.spacing.page_margin
        top = self.theme.spacing.content_top
        return Box(margin, top, self.theme.slide_width - margin * 2, self.theme.slide_height - top - 0.70)

    def add_accent_bar(self, slide: object, *, x: float | None = None, y: float | None = None, h: float = 0.62) -> None:
        deco = getattr(self.theme, "decorations", {})
        bar_w = deco.get("accent_bar_width", 0.075)
        x = 0.30 if x is None else x
        y = self.theme.spacing.title_top - 0.06 if y is None else y
        self.rect(slide, Box(x, y, bar_w, h), self.theme.colors.primary, rounded=True)
        self.rect(slide, Box(x, y + h * 0.52, bar_w, h * 0.48), self.theme.colors.accent, rounded=True)

    def add_slide_title(self, slide: object, title: str, subtitle: str = "", section: str | None = None) -> Box:
        margin = self.theme.spacing.page_margin
        title_y = self.theme.spacing.title_top
        title_font = getattr(self.theme.fonts, "title_font", "") or None
        self.add_accent_bar(slide)
        self.text(slide, Box(margin, title_y - 0.03, 9.2, 0.42), title, size=self.theme.fonts.h1_size, bold=True, font=title_font)
        if subtitle:
            self.text(slide, Box(margin, title_y + 0.45, 9.4, 0.27), subtitle, size=self.theme.fonts.subtitle_size, color=self.theme.colors.text_secondary)
        if section:
            self.add_section_label(slide, section, Box(self.theme.slide_width - margin - 0.64, title_y + 0.02, 0.64, 0.28))
        return self.content_box()

    def add_footer(self, slide: object, text: str = "SlideForge - Agent-driven PPT UI Framework") -> None:
        deco = getattr(self.theme, "decorations", {})
        margin = self.theme.spacing.page_margin
        y = self.theme.spacing.footer_y
        line_w = deco.get("footer_line_width", 0.025)
        lengths = deco.get("footer_line_lengths", [1.25, 0.72])
        self.rect(slide, Box(margin, y, lengths[0], line_w), self.theme.colors.primary, rounded=True)
        if len(lengths) > 1:
            self.rect(slide, Box(margin + lengths[0], y, lengths[1], line_w), self.theme.colors.accent, rounded=True)
        self.text(slide, Box(margin, y + 0.09, 5.2, 0.18), text, size=self.theme.fonts.tiny_size, color=self.theme.colors.text_tertiary)

    def add_card(
        self,
        slide: object,
        box: Box,
        *,
        fill: str | None = None,
        line: str | None = None,
        shadow: bool = True,
    ) -> object:
        style = self.theme.component_style("card")
        line_w = style.line_width if style.line_width is not None else 0.45
        rounded_val = True
        deco = getattr(self.theme, "decorations", {})
        if "card_radius" in deco:
            rounded_val = deco["card_radius"]

        shadow_cfg = self.theme.shadow
        if shadow and self.theme.card_shadow:
            if shadow_cfg.blur_radius <= 0:
                self.rect(
                    slide,
                    Box(box.x + shadow_cfg.card_offset_x, box.y + shadow_cfg.card_offset_y, box.w, box.h),
                    self.theme.colors.shadow_card,
                    line=None,
                    rounded=rounded_val,
                )

        shape = self.rect(
            slide, box,
            fill or style.fill,
            line or style.border,
            rounded=rounded_val,
            line_width=line_w,
            opacity=style.opacity,
            dash_style=style.dash_style,
        )

        if shadow and self.theme.card_shadow and shadow_cfg.blur_radius > 0:
            _add_native_shadow(
                shape,
                blur_radius=int(shadow_cfg.blur_radius * 914400),
                distance=int(shadow_cfg.distance * 914400),
                direction=shadow_cfg.direction,
                opacity=shadow_cfg.opacity,
            )

        top_border = deco.get("card_top_border", "")
        if top_border:
            self.rect(slide, Box(box.x, box.y, box.w, 0.04), top_border, rounded=False)

        return shape

    def add_metric_card(
        self,
        slide: object,
        box: Box,
        *,
        label: str,
        value: str,
        delta: str = "",
        note: str = "vs previous",
        icon: str = "",
        fill: str | None = None,
        border: str | None = None,
        accent: str | None = None,
    ) -> None:
        style = self.theme.component_style("metric_card")
        accent_color = accent or style.accent
        self.add_card(slide, box, fill=fill or style.fill, line=border or style.border)
        pad = self.theme.spacing.card_padding
        inner = box.inset(pad, 0.18, pad, 0.16)
        if icon:
            icon_box = Box(inner.x + inner.w - 0.34, inner.y, 0.28, 0.28)
            self.circle(slide, icon_box, self.theme.colors.primary_soft, line=self.theme.colors.border_light)
            self.text(slide, icon_box.inset(0.03, 0.02), icon, size=9, color=accent_color, bold=True, align="center", valign="middle")
        compact = box.h < 1.40
        label_size = max(8, self.theme.fonts.caption_size - (1 if compact else 0))
        caption_font = getattr(self.theme.fonts, "caption_font", "") or None
        value_size = 19 if compact else 23
        value_y = 0.34 if compact else 0.42
        delta_y = 0.78 if compact else 0.88
        note_y = 0.98 if compact else 1.12

        self.text(slide, Box(inner.x, inner.y + 0.02, inner.w - 0.45, 0.22), label, size=label_size, color=self.theme.colors.text_secondary, bold=True, font=caption_font)
        self.text(slide, Box(inner.x, inner.y + value_y, inner.w, 0.38), value, size=value_size, color=self.theme.colors.text_primary, bold=True)
        if delta:
            delta_color = self.theme.colors.success if delta.startswith("+") else self.theme.colors.warning
            delta_prefix = "+ " if delta.startswith("+") else "- " if delta.startswith("-") else ""
            delta_value = delta.lstrip("+-")
            self.text(slide, Box(inner.x, inner.y + delta_y, inner.w, 0.20), f"{delta_prefix}{delta_value}", size=9 if compact else 10, color=delta_color, bold=True)
        if note and inner.h >= 1.08:
            self.text(slide, Box(inner.x, min(inner.y + note_y, inner.y + inner.h - 0.20), inner.w, 0.18), note, size=self.theme.fonts.tiny_size, color=self.theme.colors.text_tertiary)

    def add_section_label(self, slide: object, label: str, box: Box) -> None:
        caption_font = getattr(self.theme.fonts, "caption_font", "") or None
        self.rect(slide, box, self.theme.colors.primary, rounded=True)
        self.text(slide, box.inset(0.04, 0.0), label, size=self.theme.fonts.caption_size, color="FFFFFF", bold=True, align="center", valign="middle", font=caption_font)

    def add_status_timeline_node(
        self,
        slide: object,
        box: Box,
        *,
        label: str,
        date: str,
        description: str,
        status: str = "normal",
    ) -> None:
        if status == "done":
            node_fill = self.theme.colors.primary
            card_fill = self.theme.colors.primary_soft
            line_color = self.theme.colors.primary
        elif status == "active":
            node_fill = self.theme.colors.accent
            card_fill = self.theme.colors.accent_soft
            line_color = self.theme.colors.accent
        else:
            node_fill = self.theme.colors.gray_200
            card_fill = self.theme.colors.gray_50
            line_color = self.theme.colors.border

        self.circle(slide, Box(box.x + box.w / 2 - 0.13, box.y, 0.26, 0.26), node_fill, line="FFFFFF", line_width=1.0)
        self.line(slide, box.x + box.w / 2, box.y + 0.28, box.x + box.w / 2, box.y + 0.55, color=line_color, width=0.8)
        card_h = min(max(0.75, box.h - 0.55), 1.45)
        card = Box(box.x, box.y + 0.55, box.w, card_h)
        self.add_card(slide, card, fill=card_fill, shadow=False)
        self.text(slide, Box(card.x + 0.12, card.y + 0.12, card.w - 0.24, 0.24), label, size=self.theme.fonts.body_size, bold=True, align="center")
        self.text(slide, Box(card.x + 0.12, card.y + 0.42, card.w - 0.24, 0.18), date, size=self.theme.fonts.tiny_size, color=self.theme.colors.text_tertiary, align="center")
        self.text(slide, Box(card.x + 0.14, card.y + 0.68, card.w - 0.28, 0.45), description, size=9, color=self.theme.colors.text_secondary, align="center")

    def add_process_step_card(
        self,
        slide: object,
        box: Box,
        *,
        index: int,
        title: str,
        description: str,
        output: str,
        compact: bool = False,
        output_label: str = "Output",
    ) -> None:
        self.add_card(slide, box, fill=self.theme.colors.surface_white)
        badge_size = 0.24 if compact else 0.30
        top = 0.12 if compact else 0.18
        step_badge = Box(box.x + 0.14, box.y + top, badge_size, badge_size)
        self.circle(slide, step_badge, self.theme.colors.primary_soft, line=self.theme.colors.border_light)
        self.text(slide, step_badge.inset(0.02, 0.02), f"{index:02d}", size=7 if compact else 8, color=self.theme.colors.primary, bold=True, align="center", valign="middle")
        self.text(
            slide,
            Box(box.x + 0.48, box.y + top - 0.01, box.w - 0.60, 0.20),
            title,
            size=self.theme.fonts.caption_size if compact else self.theme.fonts.body_size,
            bold=True,
        )
        desc_y = box.y + (0.46 if compact else 0.60)
        desc_h = 0.30 if compact else 0.32
        self.text(slide, Box(box.x + 0.16, desc_y, box.w - 0.32, desc_h), description, size=8 if compact else 9, color=self.theme.colors.text_secondary)
        if output:
            output_h = 0.18 if compact else 0.22
            output_box = Box(box.x + 0.16, box.y + box.h - output_h - 0.12, box.w - 0.32, output_h)
            self.rect(slide, output_box, self.theme.colors.gray_50, line=self.theme.colors.border_light, rounded=True)
            label = f"{output_label}: {output}" if output_label else output
            self.text(
                slide,
                output_box.inset(0.06, 0.02),
                label,
                size=7 if compact else 8,
                color=self.theme.colors.primary_dark,
                bold=True,
                align="center",
                valign="middle",
            )

    def add_table(self, slide: object, box: Box, headers: Sequence[str], rows: Sequence[Sequence[str]]) -> None:
        style = self.theme.component_style("table", "comparison")
        self.add_card(slide, box, fill=style.fill, line=style.border)
        inner = box.inset(0.18, 0.18)
        all_rows = [list(headers), *[list(row) for row in rows]]
        row_boxes = inner.split_rows(max(1, len(all_rows)), 0.0)
        for row_i, (row_box, values) in enumerate(zip(row_boxes, all_rows)):
            fill = style.accent if row_i == 0 else (self.theme.colors.gray_50 if row_i % 2 else self.theme.colors.surface_white)
            self.rect(slide, row_box, fill, line=self.theme.colors.border_light)
            col_boxes = row_box.split_cols(max(1, len(values)), 0.0)
            for col_i, (col_box, value) in enumerate(zip(col_boxes, values)):
                text = str(value)
                if row_i > 0 and col_i == len(values) - 1 and len(text) <= 2:
                    pill = col_box.inset(0.16, 0.10)
                    self.rect(slide, pill, self.theme.colors.accent_soft, line=self.theme.colors.accent_tint, rounded=True)
                    self.text(slide, pill, f"Recommend {text}", size=8, color=self.theme.colors.accent, bold=True, align="center", valign="middle")
                else:
                    self.text(
                        slide,
                        col_box.inset(0.06, 0.03),
                        text,
                        size=self.theme.fonts.caption_size,
                        color=style.title if row_i == 0 else self.theme.colors.text_primary,
                        bold=row_i == 0,
                        align="center",
                        valign="middle",
                    )
