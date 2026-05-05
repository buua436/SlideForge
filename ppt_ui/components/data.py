from __future__ import annotations

from dataclasses import dataclass, field

from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from ppt_ui.core.component import Component, RenderContext
from ppt_ui.core.layout import Box
from ppt_ui.core.slide import Slide


def _rgb(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _palette(ctx: RenderContext) -> list[str]:
    palette = [str(color).strip().lstrip("#") for color in getattr(ctx.theme, "chart_palette", []) if str(color).strip()]
    return palette or [ctx.theme.colors.primary, ctx.theme.colors.accent, ctx.theme.colors.success, ctx.theme.colors.warning]


@dataclass
class MetricCard(Component):
    label: str
    value: str
    delta: str = ""
    note: str = "较上期"
    icon: str = ""

    def render(self, ctx: RenderContext, box: Box) -> None:
        ctx.renderer.add_metric_card(
            ctx.slide,
            box,
            label=self.label,
            value=self.value,
            delta=self.delta,
            note=self.note,
            icon=self.icon,
        )


@dataclass
class MetricCardsSlide(Slide):
    cards: list[MetricCard] = field(default_factory=list)
    subtitle: str = ""
    scenarios: str = "实验结果摘要 / 运营指标概览 / 商业复盘 / 周报总览"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, self.subtitle)
        columns = Box(content.x, content.y + 0.12, content.w, 1.55).split_cols(max(1, len(self.cards)), ctx.theme.spacing.gutter)
        for area, card in zip(columns, self.cards):
            card.render(ctx, area)

        insight = Box(content.x, content.y + 2.08, content.w, 2.05)
        r.add_card(ctx.slide, insight, fill=ctx.theme.colors.surface_white)
        r.text(ctx.slide, Box(insight.x + 0.28, insight.y + 0.26, 3.2, 0.26), "组件使用说明", size=ctx.theme.fonts.h2_size, bold=True)
        r.text(
            ctx.slide,
            Box(insight.x + 0.28, insight.y + 0.68, insight.w - 0.56, 0.42),
            "MetricCard 用统一的 label、value、delta、compared text 和 icon placeholder 组织指标，适合快速生成信息密度稳定的数据页。",
            size=ctx.theme.fonts.body_size,
            color=ctx.theme.colors.text_secondary,
        )
        r.rect(ctx.slide, Box(insight.x + 0.28, insight.y + 1.34, insight.w - 0.56, 0.38), ctx.theme.colors.primary_soft, line=ctx.theme.colors.border_light, rounded=True)
        r.text(
            ctx.slide,
            Box(insight.x + 0.42, insight.y + 1.45, insight.w - 0.84, 0.16),
            f"适用场景：{self.scenarios}",
            size=ctx.theme.fonts.caption_size,
            color=ctx.theme.colors.primary_dark,
            bold=True,
            align="center",
            valign="middle",
        )
        r.add_footer(ctx.slide)


@dataclass
class ComparisonTableSlide(Slide):
    headers: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)
    conclusion: str = "综合功能完整性、实施周期与成本，优先推荐方案 A。"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, "用轻量表格对齐方案维度、成本、周期与推荐结果")
        table_box = Box(content.x + 0.35, content.y + 0.05, content.w - 0.70, 3.62)
        r.add_table(ctx.slide, table_box, self.headers, self.rows)
        conclusion_box = Box(content.x + 0.70, content.y + 4.08, content.w - 1.40, 0.58)
        r.add_card(ctx.slide, conclusion_box, fill=ctx.theme.colors.primary_soft, line=ctx.theme.colors.border_light, shadow=False)
        r.text(
            ctx.slide,
            conclusion_box.inset(0.18, 0.16),
            self.conclusion,
            size=ctx.theme.fonts.body_size,
            color=ctx.theme.colors.primary_dark,
            bold=True,
            align="center",
            valign="middle",
        )
        r.add_footer(ctx.slide)


@dataclass
class ChartSeries:
    name: str
    values: list[float]
    color: str = ""


@dataclass
class LineChartBlock(Component):
    categories: list[str]
    series: list[ChartSeries]
    max_value: float | None = None

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        r.add_card(ctx.slide, box, fill=ctx.theme.colors.surface_white)
        chart = box.inset(0.34, 0.34, 0.28, 0.38)
        plot = Box(chart.x + 0.45, chart.y + 0.25, chart.w - 0.65, chart.h - 0.78)
        values = [value for item in self.series for value in item.values]
        max_value = self.max_value or (max(values) if values else 1)
        min_value = min(0, min(values) if values else 0)
        span = max(max_value - min_value, 1)

        for idx in range(5):
            y = plot.y + idx * plot.h / 4
            r.line(ctx.slide, plot.x, y, plot.x + plot.w, y, ctx.theme.colors.border_light, width=0.6)
            label = int(max_value - idx * span / 4)
            r.text(ctx.slide, Box(chart.x, y - 0.06, 0.34, 0.12), str(label), size=7, color=ctx.theme.colors.text_tertiary, align="right")

        colors = _palette(ctx)
        count = max(1, len(self.categories) - 1)
        for series_idx, item in enumerate(self.series):
            color = item.color or colors[series_idx % len(colors)]
            points: list[tuple[float, float]] = []
            for idx, value in enumerate(item.values):
                x = plot.x + idx * plot.w / count
                y = plot.y + plot.h - ((value - min_value) / span) * plot.h
                points.append((x, y))
                r.circle(ctx.slide, Box(x - 0.035, y - 0.035, 0.07, 0.07), color, line="FFFFFF")
            for start, end in zip(points, points[1:]):
                r.line(ctx.slide, start[0], start[1], end[0], end[1], color, width=1.3)

        for idx, label in enumerate(self.categories):
            x = plot.x + idx * plot.w / count
            r.text(ctx.slide, Box(x - 0.18, plot.y + plot.h + 0.12, 0.36, 0.12), label, size=7, color=ctx.theme.colors.text_secondary, align="center")

        legend_x = chart.x + chart.w - 1.85
        for idx, item in enumerate(self.series):
            y = chart.y + idx * 0.18
            color = item.color or colors[idx % len(colors)]
            r.rect(ctx.slide, Box(legend_x, y + 0.04, 0.12, 0.04), color, rounded=True)
            r.text(ctx.slide, Box(legend_x + 0.17, y, 0.70, 0.12), item.name, size=7, color=ctx.theme.colors.text_secondary)


@dataclass
class LineChartSlide(Slide):
    categories: list[str] = field(default_factory=list)
    series: list[ChartSeries] = field(default_factory=list)
    subtitle: str = "适合展示趋势、实验曲线和阶段性变化"

    def render(self, ctx: RenderContext, box: Box) -> None:
        content = ctx.renderer.add_slide_title(ctx.slide, self.title, self.subtitle)
        LineChartBlock(self.categories, self.series).render(ctx, Box(content.x + 0.65, content.y + 0.05, content.w - 1.30, 4.45))
        ctx.renderer.add_footer(ctx.slide)


@dataclass
class BarChartBlock(Component):
    categories: list[str]
    series: list[ChartSeries]
    max_value: float | None = None

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        r.add_card(ctx.slide, box, fill=ctx.theme.colors.surface_white)
        chart = box.inset(0.34, 0.36, 0.28, 0.35)
        plot = Box(chart.x + 0.40, chart.y + 0.20, chart.w - 0.55, chart.h - 0.68)
        values = [value for item in self.series for value in item.values]
        max_value = self.max_value or (max(values) * 1.18 if values else 1)
        colors = _palette(ctx)

        for idx in range(5):
            y = plot.y + idx * plot.h / 4
            r.line(ctx.slide, plot.x, y, plot.x + plot.w, y, ctx.theme.colors.border_light, width=0.6)
        groups = max(1, len(self.categories))
        group_w = plot.w / groups
        bar_count = max(1, len(self.series))
        bar_w = min(0.18, group_w / (bar_count + 1.8))
        for group_idx, category in enumerate(self.categories):
            center = plot.x + group_idx * group_w + group_w / 2
            for series_idx, item in enumerate(self.series):
                value = item.values[group_idx] if group_idx < len(item.values) else 0
                h = (value / max_value) * plot.h
                x = center - (bar_count * bar_w) / 2 + series_idx * bar_w
                color = item.color or colors[series_idx % len(colors)]
                r.rect(ctx.slide, Box(x, plot.y + plot.h - h, bar_w * 0.72, h), color, rounded=True)
                r.text(ctx.slide, Box(x - 0.08, plot.y + plot.h - h - 0.18, bar_w + 0.16, 0.12), str(int(value)), size=7, color=color, bold=True, align="center")
            r.text(ctx.slide, Box(center - 0.28, plot.y + plot.h + 0.12, 0.56, 0.12), category, size=7, color=ctx.theme.colors.text_secondary, align="center")


@dataclass
class BarChartSlide(Slide):
    categories: list[str] = field(default_factory=list)
    series: list[ChartSeries] = field(default_factory=list)
    subtitle: str = "适合渠道对比、指标排行和版本差异展示"

    def render(self, ctx: RenderContext, box: Box) -> None:
        content = ctx.renderer.add_slide_title(ctx.slide, self.title, self.subtitle)
        BarChartBlock(self.categories, self.series).render(ctx, Box(content.x + 0.65, content.y + 0.05, content.w - 1.30, 4.45))
        ctx.renderer.add_footer(ctx.slide)


@dataclass
class DonutSegment:
    label: str
    value: float
    color: str = ""


@dataclass
class DonutChartSlide(Slide):
    segments: list[DonutSegment] = field(default_factory=list)
    center_label: str = "总计"
    center_value: str = ""
    subtitle: str = "适合展示占比、来源分布和结构组成"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, self.subtitle)
        panel = Box(content.x + 1.20, content.y + 0.05, content.w - 2.40, 4.45)
        r.add_card(ctx.slide, panel, fill=ctx.theme.colors.surface_white)
        chart_data = ChartData()
        chart_data.categories = [item.label for item in self.segments]
        chart_data.add_series("series", [item.value for item in self.segments])
        chart_shape = ctx.slide.shapes.add_chart(
            XL_CHART_TYPE.DOUGHNUT,
            Inches(panel.x + 0.65),
            Inches(panel.y + 0.55),
            Inches(3.0),
            Inches(3.0),
            chart_data,
        )
        chart = chart_shape.chart
        chart.has_title = False
        chart.has_legend = False
        chart.plots[0].hole_size = 62
        colors = _palette(ctx)
        for idx, point in enumerate(chart.series[0].points):
            color = self.segments[idx].color or colors[idx % len(colors)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(color)
        r.circle(ctx.slide, Box(panel.x + 1.68, panel.y + 1.57, 0.92, 0.92), "FFFFFF", line="FFFFFF")
        r.text(ctx.slide, Box(panel.x + 1.74, panel.y + 1.75, 0.80, 0.16), self.center_label, size=8, color=ctx.theme.colors.text_secondary, align="center")
        r.text(ctx.slide, Box(panel.x + 1.64, panel.y + 1.96, 1.00, 0.20), self.center_value, size=14, color=ctx.theme.colors.text_primary, bold=True, align="center")
        total = sum(item.value for item in self.segments) or 1
        for idx, item in enumerate(self.segments):
            y = panel.y + 0.72 + idx * 0.42
            color = item.color or colors[idx % len(colors)]
            r.circle(ctx.slide, Box(panel.x + 4.20, y + 0.03, 0.10, 0.10), color)
            r.text(ctx.slide, Box(panel.x + 4.42, y, 1.6, 0.16), item.label, size=9, color=ctx.theme.colors.text_secondary)
            r.text(ctx.slide, Box(panel.x + 6.10, y, 0.70, 0.16), f"{item.value / total:.1%}", size=9, color=ctx.theme.colors.text_primary, bold=True, align="right")
        r.add_footer(ctx.slide)


@dataclass
class PieChartSlide(Slide):
    segments: list[DonutSegment] = field(default_factory=list)
    subtitle: str = "适合展示占比、来源分布和结构组成"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, self.subtitle)
        panel = Box(content.x + 1.20, content.y + 0.05, content.w - 2.40, 4.45)
        r.add_card(ctx.slide, panel, fill=ctx.theme.colors.surface_white)
        chart_data = ChartData()
        chart_data.categories = [item.label for item in self.segments]
        chart_data.add_series("series", [item.value for item in self.segments])
        chart_shape = ctx.slide.shapes.add_chart(
            XL_CHART_TYPE.PIE,
            Inches(panel.x + 0.75),
            Inches(panel.y + 0.55),
            Inches(3.0),
            Inches(3.0),
            chart_data,
        )
        chart = chart_shape.chart
        chart.has_title = False
        chart.has_legend = False
        colors = _palette(ctx)
        for idx, point in enumerate(chart.series[0].points):
            color = self.segments[idx].color or colors[idx % len(colors)]
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = _rgb(color)

        total = sum(item.value for item in self.segments) or 1
        for idx, item in enumerate(self.segments):
            y = panel.y + 0.72 + idx * 0.42
            color = item.color or colors[idx % len(colors)]
            r.circle(ctx.slide, Box(panel.x + 4.20, y + 0.03, 0.10, 0.10), color)
            r.text(ctx.slide, Box(panel.x + 4.42, y, 1.6, 0.16), item.label, size=9, color=ctx.theme.colors.text_secondary)
            r.text(ctx.slide, Box(panel.x + 6.10, y, 0.70, 0.16), f"{item.value / total:.1%}", size=9, color=ctx.theme.colors.text_primary, bold=True, align="right")
        r.add_footer(ctx.slide)


@dataclass
class ProgressItem:
    label: str
    value: float
    color: str = ""


@dataclass
class ProgressBarsSlide(Slide):
    items: list[ProgressItem] = field(default_factory=list)
    subtitle: str = "适合展示项目进度、完成率和阶段健康度"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, self.subtitle)
        panel = Box(content.x + 1.45, content.y + 0.15, content.w - 2.90, 4.10)
        r.add_card(ctx.slide, panel, fill=ctx.theme.colors.surface_white)
        rows = panel.inset(0.45, 0.48, 0.45, 0.46).split_rows(max(1, len(self.items)), 0.12)
        for row, item in zip(rows, self.items):
            value = max(0, min(item.value, 1))
            r.text(ctx.slide, Box(row.x, row.y + 0.02, 1.35, 0.16), item.label, size=10, color=ctx.theme.colors.text_primary, bold=True)
            track = Box(row.x + 1.55, row.y + 0.06, row.w - 2.35, 0.10)
            r.rect(ctx.slide, track, ctx.theme.colors.gray_100, rounded=True)
            r.rect(ctx.slide, Box(track.x, track.y, track.w * value, track.h), item.color or ctx.theme.colors.primary, rounded=True)
            r.text(ctx.slide, Box(row.x + row.w - 0.55, row.y + 0.01, 0.55, 0.16), f"{int(value * 100)}%", size=9, color=ctx.theme.colors.primary_dark, bold=True, align="right")
        r.add_footer(ctx.slide)


@dataclass
class GanttTask:
    label: str
    start: int
    end: int
    color: str = ""


@dataclass
class GanttSlide(Slide):
    periods: list[str] = field(default_factory=list)
    tasks: list[GanttTask] = field(default_factory=list)
    subtitle: str = "适合展示项目计划、排期和交付节奏"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, self.subtitle)
        panel = Box(content.x + 0.90, content.y + 0.10, content.w - 1.80, 4.25)
        r.add_card(ctx.slide, panel, fill=ctx.theme.colors.surface_white)
        inner = panel.inset(0.35, 0.35, 0.35, 0.35)
        label_w = 1.30
        grid = Box(inner.x + label_w, inner.y, inner.w - label_w, inner.h)
        for idx, period in enumerate(self.periods):
            col = grid.x + idx * grid.w / max(1, len(self.periods))
            r.text(ctx.slide, Box(col, grid.y, grid.w / max(1, len(self.periods)), 0.18), period, size=8, color=ctx.theme.colors.primary_dark, bold=True, align="center")
            r.line(ctx.slide, col, grid.y + 0.28, col, grid.y + grid.h, ctx.theme.colors.border_light, width=0.6)
        rows = Box(inner.x, inner.y + 0.42, inner.w, inner.h - 0.42).split_rows(max(1, len(self.tasks)), 0.08)
        colors = _palette(ctx)
        for idx, (row, task) in enumerate(zip(rows, self.tasks)):
            r.text(ctx.slide, Box(row.x, row.y + 0.02, label_w - 0.08, 0.16), task.label, size=8, color=ctx.theme.colors.text_primary)
            r.line(ctx.slide, grid.x, row.y + row.h - 0.02, grid.x + grid.w, row.y + row.h - 0.02, ctx.theme.colors.border_light, width=0.4)
            period_count = max(1, len(self.periods))
            x = grid.x + task.start * grid.w / period_count
            w = max(0.12, (task.end - task.start) * grid.w / period_count)
            r.rect(ctx.slide, Box(x, row.y + 0.04, w, 0.12), task.color or colors[idx % len(colors)], rounded=True)
        r.add_footer(ctx.slide)


@dataclass
class HeatmapSlide(Slide):
    row_labels: list[str] = field(default_factory=list)
    col_labels: list[str] = field(default_factory=list)
    values: list[list[float]] = field(default_factory=list)
    subtitle: str = "适合展示群体、功能或场景的使用热度"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, self.subtitle)
        panel = Box(content.x + 0.95, content.y + 0.05, content.w - 1.90, 4.35)
        r.add_card(ctx.slide, panel, fill=ctx.theme.colors.surface_white)
        inner = panel.inset(0.45, 0.45, 0.45, 0.45)
        label_w = 1.05
        header_h = 0.34
        grid = Box(inner.x + label_w, inner.y + header_h, inner.w - label_w, inner.h - header_h)
        cell_w = grid.w / max(1, len(self.col_labels))
        cell_h = grid.h / max(1, len(self.row_labels))
        max_value = max([value for row in self.values for value in row], default=1)
        min_value = min([value for row in self.values for value in row], default=0)
        span = max(max_value - min_value, 1)
        for col_idx, label in enumerate(self.col_labels):
            r.text(ctx.slide, Box(grid.x + col_idx * cell_w, inner.y, cell_w, 0.18), label, size=8, color=ctx.theme.colors.text_primary, bold=True, align="center")
        for row_idx, label in enumerate(self.row_labels):
            y = grid.y + row_idx * cell_h
            r.text(ctx.slide, Box(inner.x, y + 0.06, label_w - 0.10, 0.16), label, size=8, color=ctx.theme.colors.text_primary, align="right")
            for col_idx in range(len(self.col_labels)):
                value = self.values[row_idx][col_idx] if row_idx < len(self.values) and col_idx < len(self.values[row_idx]) else 0
                ratio = (value - min_value) / span
                fill = ctx.theme.colors.primary if ratio > 0.72 else ctx.theme.colors.primary_tint if ratio > 0.38 else ctx.theme.colors.gray_100
                cell = Box(grid.x + col_idx * cell_w + 0.03, y + 0.03, cell_w - 0.06, cell_h - 0.06)
                r.rect(ctx.slide, cell, fill, line="FFFFFF", rounded=True)
                r.text(ctx.slide, cell, f"{int(value)}%", size=8, color="FFFFFF" if ratio > 0.72 else ctx.theme.colors.text_primary, bold=ratio > 0.72, align="center", valign="middle")
        r.add_footer(ctx.slide)


@dataclass
class ABComparisonSlide(Slide):
    headers: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)
    note: str = ""
    subtitle: str = "适合展示实验结果、版本差异和显著性结论"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, self.subtitle)
        table = Box(content.x + 0.95, content.y + 0.10, content.w - 1.90, 3.55)
        r.add_table(ctx.slide, table, self.headers, self.rows)
        if self.note:
            note_box = Box(content.x + 1.10, content.y + 3.95, content.w - 2.20, 0.46)
            r.add_card(ctx.slide, note_box, fill=ctx.theme.colors.success_soft, line=ctx.theme.colors.border_light, shadow=False)
            r.text(ctx.slide, note_box.inset(0.18, 0.12), self.note, size=ctx.theme.fonts.caption_size, color=ctx.theme.colors.success, bold=True, align="center", valign="middle")
        r.add_footer(ctx.slide)


@dataclass
class HighlightInsightSlide(Slide):
    summary: str = ""
    bullets: list[str] = field(default_factory=list)
    next_step: str = ""
    subtitle: str = "适合突出关键结论、原因和下一步行动"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, self.subtitle)
        panel = Box(content.x + 1.25, content.y + 0.25, content.w - 2.50, 3.85)
        r.add_card(ctx.slide, panel, fill=ctx.theme.colors.primary_soft, line=ctx.theme.colors.border)
        r.text(ctx.slide, Box(panel.x + 0.38, panel.y + 0.36, panel.w - 0.76, 0.46), self.summary, size=ctx.theme.fonts.h2_size, color=ctx.theme.colors.text_primary, bold=True)
        r.bullet_list(ctx.slide, Box(panel.x + 0.42, panel.y + 1.10, panel.w - 0.84, 1.45), self.bullets, size=ctx.theme.fonts.body_size)
        if self.next_step:
            r.rect(ctx.slide, Box(panel.x + 0.42, panel.y + 3.02, panel.w - 0.84, 0.36), "FFFFFF", line="FFFFFF", rounded=True)
            r.text(ctx.slide, Box(panel.x + 0.55, panel.y + 3.13, panel.w - 1.10, 0.14), f"下一步：{self.next_step}", size=ctx.theme.fonts.caption_size, color=ctx.theme.colors.primary, bold=True, align="center")
        r.add_footer(ctx.slide)


@dataclass
class DataAnnotationsSlide(Slide):
    annotations: list[str] = field(default_factory=list)
    note: str = ""
    subtitle: str = "适合补充数据背景、异常说明和关键事件"

    def render(self, ctx: RenderContext, box: Box) -> None:
        r = ctx.renderer
        content = r.add_slide_title(ctx.slide, self.title, self.subtitle)
        panel = Box(content.x + 1.35, content.y + 0.20, content.w - 2.70, 3.85)
        r.add_card(ctx.slide, panel, fill=ctx.theme.colors.surface_white)
        rows = panel.inset(0.42, 0.42, 0.42, 0.60).split_rows(max(1, len(self.annotations)), 0.12)
        for idx, (row, text) in enumerate(zip(rows, self.annotations), start=1):
            r.add_section_label(ctx.slide, f"{idx:02d}", Box(row.x, row.y + 0.02, 0.34, 0.22))
            r.text(ctx.slide, Box(row.x + 0.50, row.y + 0.03, row.w - 0.50, 0.18), text, size=ctx.theme.fonts.body_size, color=ctx.theme.colors.text_primary)
        if self.note:
            r.text(ctx.slide, Box(panel.x + 0.42, panel.y + panel.h - 0.42, panel.w - 0.84, 0.16), f"注：{self.note}", size=ctx.theme.fonts.tiny_size, color=ctx.theme.colors.text_tertiary)
        r.add_footer(ctx.slide)
